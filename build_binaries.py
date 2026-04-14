#!/usr/bin/env python3
"""
Build script to create executables for DocuGenius
Creates native binary for macOS and batch script for Windows
"""

import os
import sys
import subprocess
import shutil
from pathlib import Path
import tempfile

def run_command(cmd, capture_output=True):
    """Run a command and return success status"""
    try:
        if capture_output:
            result = subprocess.run(cmd, shell=True, capture_output=True, text=True)
            return result.returncode == 0, result.stdout, result.stderr
        else:
            result = subprocess.run(cmd, shell=True)
            return result.returncode == 0, "", ""
    except Exception as e:
        return False, "", str(e)

def normalize_macos_arch_name(arch_name):
    """Normalize architecture aliases used by the workflow."""
    if not arch_name:
        return None

    normalized = arch_name.strip().lower()
    aliases = {
        "x64": "x86_64",
        "amd64": "x86_64",
        "x86_64": "x86_64",
        "arm64": "arm64",
        "aarch64": "arm64",
    }
    return aliases.get(normalized, normalized)

def inspect_macos_binary_architectures(binary_path):
    """Inspect a macOS binary with `file` and return normalized architectures."""
    success, stdout, stderr = run_command(f'file -b "{binary_path}"')
    if not success:
        error_output = stderr.strip() or stdout.strip() or "unknown error"
        raise RuntimeError(f"Failed to inspect binary architecture: {error_output}")

    file_output = stdout.strip()
    architectures = []
    for arch in ("x86_64", "arm64"):
        if arch in file_output:
            architectures.append(arch)

    if not architectures:
        raise RuntimeError(f"Could not determine binary architecture from: {file_output}")

    return architectures, file_output

def create_cli_source():
    """Create the CLI source code with integrated image extraction"""
    cli_source = '''#!/usr/bin/env python3
"""
DocuGenius CLI - Document to Markdown Converter with Image Extraction
A standalone document converter for DocuGenius VS Code extension
"""

import sys
import os
import argparse
from pathlib import Path
import json
import re
import hashlib
from typing import List, Dict, Tuple, Optional

def convert_text_file(file_path):
    """Convert text-based files (just read and return content)"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return content
    except UnicodeDecodeError:
        # Try with different encodings
        for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                return content
            except UnicodeDecodeError:
                continue
        return f"# {Path(file_path).name}\\n\\nError: Could not decode file content."

def convert_json_file(file_path):
    """Convert JSON file to formatted markdown"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        markdown = f"# {Path(file_path).name}\\n\\n"
        markdown += "```json\\n"
        markdown += json.dumps(data, indent=2, ensure_ascii=False)
        markdown += "\\n```\\n"
        return markdown
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting JSON: {str(e)}"

def convert_csv_file(file_path):
    """Convert CSV file to markdown table"""
    try:
        import csv
        markdown = f"# {Path(file_path).name}\\n\\n"

        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            rows = list(reader)

        if not rows:
            return markdown + "Empty CSV file."

        # Header row
        if rows:
            markdown += "| " + " | ".join(rows[0]) + " |\\n"
            markdown += "| " + " | ".join(["---"] * len(rows[0])) + " |\\n"

            # Data rows
            for row in rows[1:]:
                # Pad row to match header length
                padded_row = row + [""] * (len(rows[0]) - len(row))
                markdown += "| " + " | ".join(padded_row[:len(rows[0])]) + " |\\n"

        return markdown
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting CSV: {str(e)}"

def convert_xml_file(file_path):
    """Convert XML file to formatted markdown"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        markdown = f"# {Path(file_path).name}\\n\\n"
        markdown += "```xml\\n"
        markdown += content
        markdown += "\\n```\\n"
        return markdown
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting XML: {str(e)}"

def convert_docx_file(file_path):
    """Convert DOCX file using python-docx"""
    try:
        from docx import Document

        doc = Document(file_path)
        file_name = Path(file_path).name

        markdown = f"# {file_name}\\n\\n"

        # Extract paragraphs
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                # Handle different paragraph styles
                style_name = paragraph.style.name.lower() if paragraph.style else ""

                if "heading 1" in style_name:
                    markdown += f"# {text}\\n\\n"
                elif "heading 2" in style_name:
                    markdown += f"## {text}\\n\\n"
                elif "heading 3" in style_name:
                    markdown += f"### {text}\\n\\n"
                elif "heading 4" in style_name:
                    markdown += f"#### {text}\\n\\n"
                elif "heading 5" in style_name:
                    markdown += f"##### {text}\\n\\n"
                elif "heading 6" in style_name:
                    markdown += f"###### {text}\\n\\n"
                else:
                    markdown += f"{text}\\n\\n"

        # Extract tables
        for table in doc.tables:
            markdown += "\\n"
            for i, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip().replace('\\n', ' ')
                    row_data.append(cell_text)

                if i == 0:
                    # Header row
                    markdown += "| " + " | ".join(row_data) + " |\\n"
                    markdown += "| " + " | ".join(["---"] * len(row_data)) + " |\\n"
                else:
                    # Data row
                    markdown += "| " + " | ".join(row_data) + " |\\n"
            markdown += "\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: python-docx library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting DOCX: {str(e)}"

def convert_excel_file(file_path):
    """Convert Excel file using openpyxl"""
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(file_path, data_only=True)
        file_name = Path(file_path).name

        markdown = f"# {file_name}\\n\\n"

        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]

            markdown += f"## {sheet_name}\\n\\n"

            # Get all rows with data
            rows = list(worksheet.iter_rows(values_only=True))
            if not rows:
                markdown += "*Empty sheet*\\n\\n"
                continue

            # Filter out completely empty rows
            non_empty_rows = []
            for row in rows:
                if any(cell is not None and str(cell).strip() for cell in row):
                    non_empty_rows.append(row)

            if not non_empty_rows:
                markdown += "*No data found*\\n\\n"
                continue

            # Find the maximum number of columns with data
            max_cols = max(len([cell for cell in row if cell is not None]) for row in non_empty_rows)

            # Create markdown table
            for i, row in enumerate(non_empty_rows):
                # Convert row to strings and pad to max_cols
                row_data = []
                for j in range(max_cols):
                    if j < len(row) and row[j] is not None:
                        row_data.append(str(row[j]).strip())
                    else:
                        row_data.append("")

                if i == 0:
                    # Header row
                    markdown += "| " + " | ".join(row_data) + " |\\n"
                    markdown += "| " + " | ".join(["---"] * len(row_data)) + " |\\n"
                else:
                    # Data row
                    markdown += "| " + " | ".join(row_data) + " |\\n"

            markdown += "\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: openpyxl library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting Excel: {str(e)}"

def convert_pptx_file(file_path):
    """Convert PowerPoint file using python-pptx"""
    try:
        from pptx import Presentation

        prs = Presentation(file_path)
        file_name = Path(file_path).name

        markdown = f"# {file_name}\\n\\n"

        for i, slide in enumerate(prs.slides, 1):
            markdown += f"## Slide {i}\\n\\n"

            # Extract text from all shapes in the slide
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                for text in slide_text:
                    # Split by lines and format appropriately
                    lines = text.split('\\n')
                    for line in lines:
                        line = line.strip()
                        if line:
                            markdown += f"{line}\\n\\n"
            else:
                markdown += "*No text content found*\\n\\n"

            markdown += "---\\n\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: python-pptx library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting PowerPoint: {str(e)}"

def convert_pdf_file(file_path):
    """Convert PDF file using pdfplumber"""
    try:
        import pdfplumber

        file_name = Path(file_path).name
        markdown = f"# {file_name}\\n\\n"

        with pdfplumber.open(file_path) as pdf:
            markdown += f"**Total Pages:** {len(pdf.pages)}\\n\\n"

            for i, page in enumerate(pdf.pages, 1):
                markdown += f"## Page {i}\\n\\n"

                try:
                    text = page.extract_text()
                    if text and text.strip():
                        # Clean up the extracted text
                        lines = text.split('\\n')
                        cleaned_lines = []
                        for line in lines:
                            line = line.strip()
                            if line:
                                cleaned_lines.append(line)

                        if cleaned_lines:
                            markdown += '\\n\\n'.join(cleaned_lines) + "\\n\\n"
                        else:
                            markdown += "*No text content found on this page*\\n\\n"
                    else:
                        markdown += "*No text content found on this page*\\n\\n"

                except Exception as page_error:
                    markdown += f"*Error extracting text from page {i}: {str(page_error)}*\\n\\n"

                markdown += "---\\n\\n"

        return markdown

    except ImportError:
        return f"# {Path(file_path).name}\\n\\nError: pdfplumber library not available"
    except Exception as e:
        return f"# {Path(file_path).name}\\n\\nError converting PDF: {str(e)}"

def extract_images_from_pdf(file_path, output_dir, min_image_size=50):
    """PDF image extraction not supported in lightweight mode"""
    return [], "PDF image extraction is not supported in lightweight mode (pdfplumber does not support image extraction)"

def convert_document_file(file_path, extract_images=True):
    """Convert document files using native Python libraries with optional image extraction"""
    file_name = Path(file_path).name
    file_ext = Path(file_path).suffix.lower()

    try:
        # First, convert the document content
        if file_ext in ['.docx']:
            content = convert_docx_file(file_path)
        elif file_ext in ['.xlsx', '.xls']:
            content = convert_excel_file(file_path)
        elif file_ext in ['.pptx']:
            content = convert_pptx_file(file_path)
        elif file_ext == '.pdf':
            content = convert_pdf_file(file_path)
        else:
            # Fallback for unsupported formats
            content = f"# {file_name}\\n\\n"
            content += f"**Document Type:** {file_ext.upper()} file\\n\\n"
            content += "This file type is not yet supported for full conversion.\\n\\n"
            content += f"- **File:** {file_name}\\n"
            content += f"- **Size:** {os.path.getsize(file_path)} bytes\\n\\n"
            return content

        # If image extraction is enabled and we have a PDF
        if extract_images and file_ext == '.pdf':
            # Note about image extraction limitation
            content += "\\n\\n<!-- Note: PDF image extraction is not supported in lightweight mode (using pdfplumber) -->\\n"

        return content

    except Exception as e:
        # Error handling - return basic info with error message
        content = f"# {file_name}\\n\\n"
        content += f"**Error converting {file_ext.upper()} file**\\n\\n"
        content += f"Error: {str(e)}\\n\\n"
        content += f"- **File:** {file_name}\\n"
        content += f"- **Size:** {os.path.getsize(file_path)} bytes\\n"
        return content


def main():
    if len(sys.argv) < 2:
        print("DocuGenius CLI - Document to Markdown Converter", file=sys.stderr)
        print("Usage: docugenius-cli <file> [extract_images]", file=sys.stderr)
        print("", file=sys.stderr)
        print("Arguments:", file=sys.stderr)
        print("  file           : Path to document file", file=sys.stderr)
        print("  extract_images : true/false to enable/disable image extraction for DOCX/PPTX/XLSX (default: true)", file=sys.stderr)
        print("", file=sys.stderr)
        print("Supported formats:", file=sys.stderr)
        print("  - Text files: .txt, .md, .markdown", file=sys.stderr)
        print("  - Data files: .json, .csv, .xml, .html", file=sys.stderr)
        print("  - Documents: .docx, .xlsx, .pptx (with image extraction), .pdf (text only)", file=sys.stderr)
        print("", file=sys.stderr)
        print("Features:", file=sys.stderr)
        print("  - Converts documents to Markdown format", file=sys.stderr)
        print("  - High-quality text extraction from PDF files (using pdfplumber)", file=sys.stderr)
        print("  - Lightweight and cross-platform consistent", file=sys.stderr)
        print("  - Fast installation and execution", file=sys.stderr)
        sys.exit(1)

    file_path = sys.argv[1]
    extract_images = True

    if len(sys.argv) > 2:
        extract_images = sys.argv[2].lower() not in ['false', 'no', '0']

    if not os.path.exists(file_path):
        print(f"Error: File not found: {file_path}", file=sys.stderr)
        sys.exit(1)

    file_ext = Path(file_path).suffix.lower()

    try:
        if file_ext in ['.txt', '.md', '.markdown']:
            content = convert_text_file(file_path)
        elif file_ext == '.json':
            content = convert_json_file(file_path)
        elif file_ext == '.csv':
            content = convert_csv_file(file_path)
        elif file_ext in ['.xml', '.html', '.htm']:
            content = convert_xml_file(file_path)
        elif file_ext in ['.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.pdf']:
            content = convert_document_file(file_path, extract_images)
        else:
            # Default to text file handling for unknown extensions
            content = convert_text_file(file_path)

        print(content)

    except Exception as e:
        print(f"Error processing file: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
'''
    return cli_source

def create_darwin_binary(expected_arch=None):
    """Create macOS binary using PyInstaller"""
    print("Building DocuGenius macOS binary...")

    # Detect current architecture
    import platform
    current_arch = platform.machine()
    print(f"Current system architecture: {current_arch}")
    expected_arch = normalize_macos_arch_name(expected_arch)
    if expected_arch:
        print(f"Expected output architecture: {expected_arch}")

    # Use the shared converter.py as entry point for PyInstaller
    cli_file = str(Path('bin/converter.py').resolve())

    try:
        # Create virtual environment for building
        env_dir = "build_env_darwin"
        if os.path.exists(env_dir):
            shutil.rmtree(env_dir)

        python_cmd = sys.executable
        print(f"Using Python: {python_cmd}")
        print(f"Creating build environment: {env_dir}")
        success, _, _ = run_command(f'"{python_cmd}" -m venv {env_dir}')
        if not success:
            print("Failed to create virtual environment")
            return False

        # Install PyInstaller and document processing libraries
        print("Installing PyInstaller and document libraries...")
        install_cmd = f". {env_dir}/bin/activate && pip install pyinstaller python-docx python-pptx openpyxl pdfplumber"
        success, _, _ = run_command(install_cmd)

        if not success:
            print("Failed to install required libraries")
            return False

        # Build the executable
        print("Building executable...")
        build_cmd = f". {env_dir}/bin/activate && python -m PyInstaller --onefile --name docugenius-cli --strip --optimize=2 {cli_file}"
        success, stdout, stderr = run_command(build_cmd, capture_output=True)

        if not success:
            print("Failed to build executable")
            if stdout:
                print("STDOUT:", stdout)
            if stderr:
                print("STDERR:", stderr)
            return False

        # Check if the executable was created
        exe_path = "dist/docugenius-cli"

        if not os.path.exists(exe_path):
            print("Executable not found after build")
            return False

        # Create the bin/darwin directory if it doesn't exist
        darwin_dir = Path("bin/darwin")
        darwin_dir.mkdir(parents=True, exist_ok=True)

        # Copy the executable to the bin directory
        target_path = darwin_dir / "docugenius-cli"

        shutil.copy2(exe_path, target_path)
        os.chmod(target_path, 0o755)

        built_architectures, file_output = inspect_macos_binary_architectures(target_path)
        print(f"Binary architecture: {', '.join(built_architectures)}")
        print(f"`file` output: {file_output}")

        if expected_arch and built_architectures != [expected_arch]:
            raise RuntimeError(
                f"Expected a {expected_arch} binary, but built {', '.join(built_architectures)}. "
                "Check the GitHub Actions runner label and Python architecture."
            )

        print(f"Binary created: {target_path}")
        print(f"File size: {os.path.getsize(target_path) / (1024*1024):.1f} MB")

        # Clean up build artifacts
        cleanup_dirs = ['build', 'dist', env_dir]
        for dir_name in cleanup_dirs:
            if os.path.exists(dir_name):
                shutil.rmtree(dir_name)

        for spec_file in Path('.').glob('*.spec'):
            spec_file.unlink()

        print("Cleaned up build artifacts")
        return True

    except Exception as e:
        print(f"Failed to create binary: {e}")
        return False

def create_windows_batch():
    """Create Windows batch file"""
    print("Creating DocuGenius Windows batch file...")

    # Create the bin/win32 directory if it doesn't exist
    win32_dir = Path("bin/win32")
    win32_dir.mkdir(parents=True, exist_ok=True)

    batch_content = '''@echo off
REM DocuGenius CLI for Windows
chcp 65001 >nul 2>&1
setlocal

set PYTHONIOENCODING=utf-8
set PYTHONUTF8=1

if "%~1"=="" (
    echo DocuGenius CLI - Document to Markdown Converter with Image Extraction
    echo Usage: docugenius-cli ^<file^> [extract_images] [output_path]
    echo.
    echo Supported formats:
    echo   - Text files: .txt, .md, .markdown
    echo   - Data files: .json, .csv, .xml, .html
    echo   - Documents: .docx, .xlsx, .pptx, .pdf
    echo.
    echo Note: Python must be installed and available in PATH.
    exit /b 1
)

if not exist "%~1" (
    echo Error: File not found: "%~1"
    exit /b 1
)

python --version >nul 2>&1
if errorlevel 1 (
    echo Error: Python is not installed or not in PATH
    echo Please install Python from https://python.org
    exit /b 1
)

set "SCRIPT_DIR=%~dp0"
set "CONVERTER=%SCRIPT_DIR%..\\converter.py"

if not exist "%CONVERTER%" (
    echo Error: converter.py not found at "%CONVERTER%"
    echo Please reinstall the DocuGenius extension.
    exit /b 1
)

python "%CONVERTER%" "%~1" "%~2" "%~3"
exit /b %ERRORLEVEL%
'''

    target_path = win32_dir / "docugenius-cli.bat"

    try:
        with open(target_path, 'w') as f:
            f.write(batch_content)

        print(f"Windows batch file created: {target_path}")
        print(f"File size: {os.path.getsize(target_path)} bytes")
        return True

    except Exception as e:
        print(f"Failed to create batch file: {e}")
        return False

def main():
    print("DocuGenius Binary Builder")

    current_platform = sys.platform

    if len(sys.argv) > 1:
        target = sys.argv[1].lower()
    else:
        target = "all"

    expected_arch = normalize_macos_arch_name(sys.argv[2]) if len(sys.argv) > 2 else None

    success = True

    if target in ["all", "darwin", "macos"]:
        if current_platform == "darwin" or target != "all":
            success &= create_darwin_binary(expected_arch=expected_arch)
        else:
            print("Skipping macOS binary (not on macOS)")

    if target in ["all", "windows", "win32"]:
        success &= create_windows_batch()

    if success:
        print("\nBinary build completed successfully!")
        print("\nGenerated files:")
        if os.path.exists("bin/darwin/docugenius-cli"):
            size = os.path.getsize("bin/darwin/docugenius-cli") / (1024*1024)
            print(f"   - bin/darwin/docugenius-cli ({size:.1f} MB)")
        if os.path.exists("bin/win32/docugenius-cli.bat"):
            size = os.path.getsize("bin/win32/docugenius-cli.bat")
            print(f"   - bin/win32/docugenius-cli.bat ({size} bytes)")
    else:
        print("\nBuild failed!")
        sys.exit(1)

if __name__ == "__main__":
    main()
